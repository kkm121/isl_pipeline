"""
=============================================================================
BharatSRM-Net v4: SIH 2026 Presentation Generator (Team: ChaosToCode)
=============================================================================
Updates: C:\\Users\\muthu\\Downloads\\SIH2026-IDEA-Presentation-Format (2).pptx
Problem Statement: NTRO PS ID 26142 (Deep Learning Super-Resolution Mapping)
=============================================================================
"""

import os
import sys
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPT_PATH = r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx"

def style_text_frame(tf, items, font_size=13, space_after=8, text_color=RGBColor(30, 41, 59), bold_color=RGBColor(15, 23, 42)):
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
        p.space_after = Pt(space_after)
        p.font.size = Pt(font_size)
        p.font.name = "Segoe UI"
        
        if ":" in item:
            parts = item.split(":", 1)
            r1 = p.add_run()
            r1.text = parts[0] + ":"
            r1.font.bold = True
            r1.font.size = Pt(font_size)
            r1.font.name = "Segoe UI"
            r1.font.color.rgb = bold_color
            
            r2 = p.add_run()
            r2.text = parts[1]
            r2.font.bold = False
            r2.font.size = Pt(font_size)
            r2.font.name = "Segoe UI"
            r2.font.color.rgb = text_color
        else:
            r = p.add_run()
            r.text = item
            r.font.bold = False
            r.font.size = Pt(font_size)
            r.font.name = "Segoe UI"
            r.font.color.rgb = text_color

def update_presentation():
    if not os.path.exists(PPT_PATH):
        print(f"Error: {PPT_PATH} not found!")
        return

    prs = pptx.Presentation(PPT_PATH)
    print(f"Loaded presentation with {len(prs.slides)} slides.")

    # -------------------------------------------------------------
    # SLIDE 1: Title Page
    # -------------------------------------------------------------
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if shape.has_text_frame:
            if "TITLE PAGE" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "BharatSRM-Net v4"
                p.font.size = Pt(28)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
                
                p2 = shape.text_frame.add_paragraph()
                p2.text = "Physically-Consistent, Uncertainty-Aware Super-Resolution Framework for Indian Satellite Imagery"
                p2.font.size = Pt(14)
                p2.font.name = "Segoe UI"
                p2.font.color.rgb = RGBColor(71, 85, 105)
                
            elif "SMART INDIA HACKATHON" in shape.text:
                p = shape.text_frame.paragraphs[0]
                p.text = "SMART INDIA HACKATHON 2026"
                p.font.size = Pt(22)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)
                
            elif "Problem Statement ID" in shape.text:
                s1_items = [
                    "Problem Statement ID: 26142",
                    "Problem Statement Title: Deep Learning Based Super Resolution Mapping (SRM) from Medium Resolution Satellite Imageries",
                    "Theme: Smart Education / Space Technology & Defense",
                    "PS Category: Software",
                    "Organization / Ministry: National Technical Research Organisation (NTRO)",
                    "Team Name: ChaosToCode",
                    "Team ID: [To be filled by team]",
                ]
                style_text_frame(shape.text_frame, s1_items, font_size=13, space_after=6)

    # -------------------------------------------------------------
    # SLIDE 2: Proposed Solution
    # -------------------------------------------------------------
    slide2 = prs.slides[1]
    for shape in slide2.shapes:
        if shape.has_text_frame:
            if "IDEA TITLE" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "PROPOSED SOLUTION: BHARATSRM-NET v4"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif "Proposed Solution" in shape.text:
                s2_items = [
                    "4x Physical Super-Resolution (10m -> 2.5m GSD): Enhances freely accessible Sentinel-2 multispectral imagery to 2.5m GSD (<4m target requirement), providing 16x higher spatial sampling density.",
                    "Zero-DC Residual Learning: Decouples spatial luminance (Y) from spectral chrominance (Cr, Cb) to guarantee 100% true-color fidelity with mathematically zero spectral drift or false color casts.",
                    "Sensor MTF Degradation Consistency (L_degrade): Enforces Point Spread Function (PSF) and sensor MTF invariance, ensuring super-resolved imagery physically degrades back to low-resolution sensor physics.",
                    "Calibrated Heteroscedastic Uncertainty (sigma^2): Quantifies per-pixel error variance via temperature-calibrated negative log-likelihood (Kendall & Gal), proving trust and preventing AI hallucinations for defense analysts.",
                    "Multi-Task Downstream Value: Integrated neural heads for PMGSY Rural Road Extraction, ISRO 5-Class LULC Disaggregation, and Disaster Change Assessment running directly from shared 64-dim backbone features.",
                ]
                style_text_frame(shape.text_frame, s2_items, font_size=12, space_after=6)
            elif "Your Team Name" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "ChaosToCode"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)

    # -------------------------------------------------------------
    # SLIDE 3: Technical Approach
    # -------------------------------------------------------------
    slide3 = prs.slides[2]
    for shape in slide3.shapes:
        if shape.has_text_frame:
            if "TECHNICAL APPROACH" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "TECHNICAL APPROACH & SYSTEM ARCHITECTURE"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif "Technologies to be used" in shape.text:
                s3_items = [
                    "4-Plane Swarm Architecture: Control Plane (Multi-Agent Swarm), Integration Plane (Cloud-Optimized GeoTIFF Streaming), Execution Plane (PyTorch, CUDA AMP FP16), Verification Plane (Deterministic Gate Enforcement).",
                    "Neural Network Backbone: Masked PartialConv2d (QA60/S2cloudless cloud resilience) -> AC-FEM Cross-Attention (CartoDEM elevation & slope fusion) -> Dilated Residual Blocks (r in {1,2,4,8}) -> ICNR PixelShuffle (s=4).",
                    "5-Term Composite Loss Function: L_total = 1.0 L_rec (Charbonnier) + 0.2 L_spec (SAM) + 0.3 L_degrade (Sensor PSF) + 0.5 L_struct (Laplacian Edge) + 0.01 L_conf (NLL Uncertainty).",
                    "Tiled 2D Hanning Inference Engine: Windowed tiling with 2D Hanning reconstruction enabling seamless inference over 10,000x10,000 regional scenes with zero boundary checkerboard seams.",
                    "Interactive Web GIS Studio: Full-stack FastAPI backend + interactive Split-Screen GIS Studio with instant GeoTIFF and GeoJSON vector export capabilities.",
                ]
                style_text_frame(shape.text_frame, s3_items, font_size=12, space_after=6)
            elif "Your Team Name" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "ChaosToCode"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)

    # -------------------------------------------------------------
    # SLIDE 4: Feasibility and Viability
    # -------------------------------------------------------------
    slide4 = prs.slides[3]
    for shape in slide4.shapes:
        if shape.has_text_frame:
            if "FEASIBILITY AND VIABILITY" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "FEASIBILITY, VIABILITY & RISK MITIGATION"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif "Analysis of the feasibility" in shape.text:
                s4_items = [
                    "Technical Feasibility (Verified Working Prototype): Complete end-to-end framework fully built and tested; sub-second GPU inference per scene verified across 4 diverse Indian geographic biomes.",
                    "Data Sourcing & Open Pipeline: Ingests 100% free, publicly accessible Copernicus Sentinel-2 L2A BOA data and ISRO Bhuvan CartoDEM tiles, requiring zero proprietary API licensing costs.",
                    "Risk 1 - AI Hallucinations in Defense: Solved via Calibrated Uncertainty Quantification (sigma^2) and Sensor PSF cycle-consistency, alerting defense analysts to sub-pixel boundary ambiguity.",
                    "Risk 2 - Cloud & Atmospheric Occlusion: Solved via PartialConv2d invalid-pixel masking combined with CartoDEM topographic elevation priors to reconstruct underlying terrain contours.",
                    "Risk 3 - Compute & Latency on Regional Scales: Solved via INT8 dynamic quantization and windowed 2D Hanning tiled streaming engines.",
                    "Deployment Viability: Sealed Docker containerization (--network=none --read-only) ensuring seamless deployment on air-gapped defense servers, cloud platforms, and edge drone hardware.",
                ]
                style_text_frame(shape.text_frame, s4_items, font_size=12, space_after=5)
            elif "Your Team Name" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "ChaosToCode"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)

    # -------------------------------------------------------------
    # SLIDE 5: Impact and Benefits
    # -------------------------------------------------------------
    slide5 = prs.slides[4]
    for shape in slide5.shapes:
        if shape.has_text_frame:
            if "IMPACT AND BENEFITS" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "IMPACT, BENEFITS & NATIONAL VALUE"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif "Potential impact" in shape.text:
                s5_items = [
                    "Strategic Defense Impact (NTRO): Delivers sovereign sub-4m (2.5m) high-resolution intelligence from free open data, saving crores of rupees in commercial satellite procurement (> $15/km^2 for foreign imagery).",
                    "Anti-Hallucination Assurance: Mathematical confidence heatmaps prevent false military target identification, enabling selective human-in-the-loop verification on only high-uncertainty zones.",
                    "PMGSY Rural Road Infrastructure: Automated extraction and vectorization of unpaved village roads and connectivity corridors for PM Gati Shakti National Master Plan.",
                    "Agriculture & Food Security: Precision crop parcel boundary delineation and individual field acreage monitoring for PM Fasal Bima Yojana.",
                    "Disaster Management & Flood Response: Rapid flood inundation extent mapping, shoreline retreat tracking, and infrastructure damage triage during extreme weather events.",
                    "ISRO LULC & Urban Planning: Automated tracking of urban sprawl, water reservoir depletion, and forest canopy conservation.",
                ]
                style_text_frame(shape.text_frame, s5_items, font_size=12, space_after=5)
            elif "Your Team Name" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "ChaosToCode"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)

    # -------------------------------------------------------------
    # SLIDE 6: Research and References
    # -------------------------------------------------------------
    slide6 = prs.slides[5]
    for shape in slide6.shapes:
        if shape.has_text_frame:
            if "RESEARCH" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "RESEARCH, DATASETS & REFERENCES"
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(15, 23, 42)
            elif "Details / Links" in shape.text:
                s6_items = [
                    "Bayesian Uncertainty Estimation: Kendall & Gal (NeurIPS) - 'What Uncertainties Do We Need in Bayesian Deep Learning for Computer Vision?'",
                    "Irregular Hole & Cloud Inpainting: Liu et al. (ECCV) - 'Image Inpainting for Irregular Holes Using Partial Convolutions.'",
                    "Sub-Pixel Anti-Aliasing: Odena et al. - 'Deconvolution and Checkerboard Artifacts' (ICNR Sub-Pixel Convolution Formulation).",
                    "Earth Observation Super-Resolution Benchmark: Cornebise et al. (NeurIPS) - 'WorldStrat: A Dataset for Spatial Super-Resolution in Earth Observation.'",
                    "Operational Data Sources & Portals: Copernicus Open Access Hub (Sentinel-2 Level-2A BOA Reflectance), ISRO NRSC Bhuvan (CartoDEM & Indian LULC Classification Standards), SPOT 6/7 1.5m Pansharpened RGBN References.",
                ]
                style_text_frame(shape.text_frame, s6_items, font_size=12, space_after=6)
            elif "Your Team Name" in shape.text:
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                p.text = "ChaosToCode"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.name = "Segoe UI"
                p.font.color.rgb = RGBColor(2, 132, 199)

    # Save modified presentation
    prs.save(PPT_PATH)
    print(f"[SUCCESS] Updated and saved presentation to: {PPT_PATH}")

if __name__ == "__main__":
    update_presentation()
