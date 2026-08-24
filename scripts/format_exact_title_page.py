"""
=============================================================================
Update Slide 1: Exact 18pt font size across all title page fields
=============================================================================
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_FILES = [
    r"C:\Users\muthu\Downloads\SIH2026-Presentation-Final-18pt.pptx",
    r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-ChaosToCode.pptx",
    r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx",
    r"C:\Users\muthu\Downloads\SIH2026-BharatSRM-Net-PitchDeck.pptx",
    r"C:\Users\muthu\Downloads\SIH2026-BharatSRM-Net-SmartEducation.pptx",
]

def format_title_page_18pt(out_path):
    template_path = r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx"
    if not os.path.exists(template_path):
        return
    try:
        prs = pptx.Presentation(template_path)
        s1 = prs.slides[0]
        
        for shape in s1.shapes:
            if shape.has_text_frame:
                if "TITLE PAGE" in shape.text or "SMART INDIA" in shape.text or "BharatSRM" in shape.text:
                    if shape.name == "Subtitle 3":
                        shape.text_frame.clear()
                    else:
                        shape.top = Inches(0.4)
                        shape.left = Inches(0.5)
                        shape.width = Inches(8.0)
                        shape.height = Inches(0.7)
                        shape.text_frame.clear()
                        p = shape.text_frame.paragraphs[0]
                        p.text = "TITLE PAGE"
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.name = "Arial"
                        p.font.color.rgb = RGBColor(0, 0, 0)
                elif "Problem Statement ID" in shape.text:
                    shape.top = Inches(1.3)
                    shape.left = Inches(0.5)
                    shape.width = Inches(12.0)
                    shape.height = Inches(5.5)
                    tf = shape.text_frame
                    tf.clear()
                    
                    fields = [
                        ("Problem Statement ID \u2013 ", "26142"),
                        ("Problem Statement Title- ", "Deep Learning Based Super Resolution Mapping (SRM) from Medium Resolution Satellite Imageries"),
                        ("Theme- ", "Smart Education"),
                        ("PS Category- ", "Software"),
                        ("Team ID- ", "1"),
                        ("Team Name (Registered on portal)- ", "CHAOS TO CODE"),
                    ]
                    
                    for idx, (label, val) in enumerate(fields):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        p.space_after = Pt(14)
                        p.font.name = "Arial"
                        
                        r1 = p.add_run()
                        r1.text = label
                        r1.font.bold = True
                        r1.font.size = Pt(18) # 18pt font size
                        r1.font.name = "Arial"
                        r1.font.color.rgb = RGBColor(0, 0, 0)
                        
                        r2 = p.add_run()
                        r2.text = val
                        r2.font.bold = True
                        r2.font.size = Pt(18) # 18pt font size
                        r2.font.name = "Arial"
                        r2.font.color.rgb = RGBColor(16, 56, 107) # SIH Navy Blue

        prs.save(out_path)
        print(f"[SUCCESS] Updated Title Page (18pt) for: {out_path}")
    except Exception as e:
        print(f"[NOTE] Could not update {out_path}: {e}")

if __name__ == "__main__":
    for p in PPT_FILES:
        format_title_page_18pt(p)
