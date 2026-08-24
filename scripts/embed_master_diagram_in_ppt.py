"""
=============================================================================
Embed Master Architecture Diagram into SIH Presentation
=============================================================================
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-ChaosToCode.pptx"
ALT_PPT_PATH = r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx"
MASTER_DIAG_PATH = r"outputs/master_architecture_diagram.png"

def embed_master_architecture():
    for path in [PPT_PATH, ALT_PPT_PATH]:
        if not os.path.exists(path):
            continue
        try:
            prs = pptx.Presentation(path)
            slide3 = prs.slides[2] # Slide 3: Technical Approach
            
            # Remove any previous diagram images on Slide 3
            for s in list(slide3.shapes):
                if s.name.startswith("Added_Diagram"):
                    slide3.shapes._spTree.remove(s._element)

            # Adjust title and text box to leave maximum space for the master architecture diagram
            for shape in slide3.shapes:
                if shape.has_text_frame:
                    if "TECHNICAL APPROACH" in shape.text:
                        shape.top = Inches(0.15)
                        shape.height = Inches(0.7)
                    elif "Multi-Modal Ingestion" in shape.text or "Technologies to be used" in shape.text or "4-Plane" in shape.text:
                        shape.top = Inches(0.85)
                        shape.height = Inches(1.5)
                        shape.width = Inches(12.5)
                        shape.left = Inches(0.4)

            # Add High-Resolution Master Architecture Diagram
            if os.path.exists(MASTER_DIAG_PATH):
                pic = slide3.shapes.add_picture(
                    MASTER_DIAG_PATH,
                    Inches(0.4),
                    Inches(2.45),
                    width=Inches(12.5),
                    height=Inches(4.45)
                )
                pic.name = "Added_Diagram_Master_Architecture"
                print(f"[OK] Added Master Architecture Diagram to Slide 3 in {path}")

            prs.save(path)
            print(f"[SUCCESS] Saved updated deck: {path}")
        except Exception as e:
            print(f"[NOTE] Error on {path}: {e}")

if __name__ == "__main__":
    embed_master_architecture()
