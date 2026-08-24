"""
=============================================================================
BharatSRM-Net v4: Pixel-Perfect 2-Column Presentation Layout (Zero Overlap)
=============================================================================
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_FILES = [
    r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-ChaosToCode.pptx",
    r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx",
]
DIAG_DIR = "outputs/ppt_diagrams"
MASTER_DIAG = "outputs/master_architecture_diagram.png"

def style_bullet(p, title, desc, font_size=10.5, space_after=4, bold_color=RGBColor(15, 23, 42), text_color=RGBColor(51, 65, 85)):
    p.text = ""
    p.space_after = Pt(space_after)
    p.font.name = "Segoe UI"
    
    r1 = p.add_run()
    r1.text = f"• {title}: "
    r1.font.bold = True
    r1.font.size = Pt(font_size)
    r1.font.name = "Segoe UI"
    r1.font.color.rgb = bold_color
    
    r2 = p.add_run()
    r2.text = desc
    r2.font.bold = False
    r2.font.size = Pt(font_size)
    r2.font.name = "Segoe UI"
    r2.font.color.rgb = text_color

def apply_clean_layout(path):
    if not os.path.exists(path):
        return
    try:
        prs = pptx.Presentation(path)
        print(f"\nProcessing {path}...")

        # -------------------------------------------------------------
        # SLIDE 1: Title Slide
        # -------------------------------------------------------------
        s1 = prs.slides[0]
        for shape in s1.shapes:
            if shape.has_text_frame:
                if "SMART INDIA HACKATHON" in shape.text:
                    shape.top = Inches(0.35)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.7)
                    shape.text_frame.paragraphs[0].font.size = Pt(22)
                elif "TITLE PAGE" in shape.text or "BharatSRM" in shape.text:
                    shape.top = Inches(0.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(9.0)
                    shape.height = Inches(0.9)
                    shape.text_frame.clear()
                    p1 = shape.text_frame.paragraphs[0]
                    p1.text = "BharatSRM-Net v4"
                    p1.font.size = Pt(26)
                    p1.font.bold = True
                    p1.font.color.rgb = RGBColor(15, 23, 42)
                    p2 = shape.text_frame.add_paragraph()
                    p2.text = "Physically-Consistent, Uncertainty-Aware Super-Resolution Framework for Indian Satellite Imagery"
                    p2.font.size = Pt(12.5)
                    p2.font.color.rgb = RGBColor(71, 85, 105)
                elif "Problem Statement ID" in shape.text:
                    shape.top = Inches(1.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(6.8)
                    shape.height = Inches(4.5)
                    tf = shape.text_frame
                    tf.clear()
                    fields = [
                        ("Problem Statement ID", "26142"),
                        ("Problem Statement Title", "Deep Learning Based Super Resolution Mapping (SRM) from Medium Resolution Satellite Imageries"),
                        ("Theme", "Smart Education / Space Technology & Defense"),
                        ("PS Category", "Software"),
                        ("Organization / Ministry", "National Technical Research Organisation (NTRO)"),
                        ("Team Name", "ChaosToCode"),
                        ("Team ID", "[To be filled by team]"),
                    ]
                    for idx, (k, v) in enumerate(fields):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=11, space_after=3, bold_color=RGBColor(2, 132, 199))

        # -------------------------------------------------------------
        # SLIDE 2: Proposed Solution (2-Column Layout)
        # -------------------------------------------------------------
        s2 = prs.slides[1]
        for s in list(s2.shapes):
            if s.name.startswith("Added_Diagram"):
                s2.shapes._spTree.remove(s._element)

        for shape in s2.shapes:
            if shape.has_text_frame:
                if "PROPOSED SOLUTION" in shape.text or "IDEA TITLE" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                elif "The Problem" in shape.text or "Proposed Solution" in shape.text or "4x Physical" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(5.3)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("The Core Challenge", "Free 10m Sentinel-2 satellite imagery is too blurry for narrow village roads, farm parcel boundaries, and defense observation, while commercial 2.5m imagery costs crores (> ₹12 Lakh per 1,000 km²)."),
                        ("The Breakthrough", "BharatSRM-Net v4 transforms free 10m Sentinel-2 data into 2.5m commercial-grade imagery (<4m NTRO target requirement) with 16x higher spatial sampling density and mathematically 0% color drift."),
                        ("Defense Anti-Hallucination", "Outputs a Calibrated Uncertainty Heatmap (σ²) proving to intelligence analysts exactly which reconstructed pixels are 100% trustworthy."),
                        ("Multi-Task AI Value", "Integrated neural heads for PMGSY Rural Roads, ISRO 5-Class LULC Disaggregation, and Disaster Flood Mapping running from the same backbone."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=5)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag2 = f"{DIAG_DIR}/diagram_slide2.png"
        if os.path.exists(diag2):
            pic = s2.shapes.add_picture(diag2, Inches(6.0), Inches(1.45), width=Inches(6.8), height=Inches(4.8))
            pic.name = "Added_Diagram_2"

        # -------------------------------------------------------------
        # SLIDE 3: Technical Approach & Master Architecture
        # -------------------------------------------------------------
        s3 = prs.slides[2]
        for s in list(s3.shapes):
            if s.name.startswith("Added_Diagram"):
                s3.shapes._spTree.remove(s._element)

        for shape in s3.shapes:
            if shape.has_text_frame:
                if "TECHNICAL APPROACH" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.7)
                elif "Multi-Modal Ingestion" in shape.text or "Technologies to be used" in shape.text or "4-Plane" in shape.text:
                    shape.top = Inches(0.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(12.3)
                    shape.height = Inches(1.0)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Multi-Modal Ingestion & AI Core", "10-Band Sentinel-2 BOA + CartoDEM terrain slope/aspect priors -> PartialConv2d cloud masking -> AC-FEM Cross-Attention -> Dilated ResBlocks -> ICNR PixelShuffle (s=4)."),
                        ("Production Deployment", "Tiled 2D Hanning inference engine for 10,000x10,000 scenes with zero tile boundary seams + Interactive Web GIS Split-Screen Studio."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=2)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        if os.path.exists(MASTER_DIAG):
            pic = s3.shapes.add_picture(MASTER_DIAG, Inches(0.5), Inches(2.1), width=Inches(12.3), height=Inches(4.6))
            pic.name = "Added_Diagram_Master_Architecture"

        # -------------------------------------------------------------
        # SLIDE 4: Feasibility, Viability & Deployment (2-Column Layout)
        # -------------------------------------------------------------
        s4 = prs.slides[3]
        for s in list(s4.shapes):
            if s.name.startswith("Added_Diagram"):
                s4.shapes._spTree.remove(s._element)

        for shape in s4.shapes:
            if shape.has_text_frame:
                if "FEASIBILITY" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                elif "Verified Working" in shape.text or "Analysis of the feasibility" in shape.text or "Technical Feasibility" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(5.3)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Verified Working Prototype", "Full-stack framework built and verified with live sub-second GPU inference across 4 Indian biomes."),
                        ("100% Free Data Pipeline", "Uses openly accessible Copernicus Sentinel-2 L2A BOA and ISRO Bhuvan CartoDEM data (zero API licensing cost)."),
                        ("Risk 1 - Defense AI Hallucination", "Solved via Calibrated Uncertainty Quantification (σ²) alerting analysts to boundary ambiguity."),
                        ("Risk 2 - Cloud Occlusion", "Solved via PartialConv2d cloud masking + CartoDEM topographic elevation context."),
                        ("Deployment Viability", "Sealed Docker containerization (--network=none --read-only) ensuring zero data leakage on air-gapped defense servers."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10, space_after=4)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag4 = f"{DIAG_DIR}/diagram_slide4.png"
        if os.path.exists(diag4):
            pic = s4.shapes.add_picture(diag4, Inches(6.0), Inches(1.45), width=Inches(6.8), height=Inches(4.8))
            pic.name = "Added_Diagram_4"

        # -------------------------------------------------------------
        # SLIDE 5: Business Impact & Competitive ROI Table (2-Column Layout)
        # -------------------------------------------------------------
        s5 = prs.slides[4]
        for s in list(s5.shapes):
            if s.name.startswith("Added_Diagram"):
                s5.shapes._spTree.remove(s._element)

        for shape in s5.shapes:
            if shape.has_text_frame:
                if "IMPACT AND BENEFITS" in shape.text or "BUSINESS IMPACT" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                elif "Massive Cost Savings" in shape.text or "Potential impact" in shape.text or "Strategic Defense" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(4.9)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Massive Cost Savings (ROI)", "Saves ₹100s of Crores annually by replacing expensive commercial satellite purchases ($15/km²)."),
                        ("PMGSY Rural Infrastructure", "Automated vectorization of unpaved village roads and connectivity corridors for PM Gati Shakti."),
                        ("Agriculture & Food Security", "Precision crop parcel boundary tracking (PM Fasal Bima) + individual field acreage monitoring."),
                        ("Disaster Response", "Rapid flood inundation extent mapping and infrastructure damage triage during extreme weather events."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=5)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag5 = f"{DIAG_DIR}/diagram_slide5.png"
        if os.path.exists(diag5):
            pic = s5.shapes.add_picture(diag5, Inches(5.6), Inches(1.4), width=Inches(7.2), height=Inches(5.0))
            pic.name = "Added_Diagram_5"

        # -------------------------------------------------------------
        # SLIDE 6: Research & References
        # -------------------------------------------------------------
        s6 = prs.slides[5]
        for shape in s6.shapes:
            if shape.has_text_frame:
                if "RESEARCH" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                elif "Bayesian Deep Learning" in shape.text or "Details / Links" in shape.text or "Bayesian Uncertainty" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.6)
                    shape.width = Inches(12.0)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Bayesian Uncertainty Estimation", "Kendall & Gal (NeurIPS) - 'What Uncertainties Do We Need in Bayesian Deep Learning for Computer Vision?'"),
                        ("Irregular Cloud & Hole Inpainting", "Liu et al. (ECCV) - 'Image Inpainting for Irregular Holes Using Partial Convolutions.'"),
                        ("Sub-Pixel Anti-Aliasing", "Odena et al. - 'Deconvolution and Checkerboard Artifacts' (ICNR Sub-Pixel Convolution Formulation)."),
                        ("Earth Observation Super-Resolution Benchmark", "Cornebise et al. (NeurIPS) - 'WorldStrat: A Dataset for Spatial Super-Resolution in Earth Observation.'"),
                        ("Operational Portals & Ground Truth", "Copernicus Open Access Hub (Sentinel-2 BOA), ISRO NRSC Bhuvan (CartoDEM & LULC Schema), SPOT 6/7 1.5m Pansharpened RGBN."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=11.5, space_after=6)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        prs.save(path)
        print(f"[SUCCESS] Cleaned layout with zero overlap saved to: {path}")
    except Exception as e:
        print(f"[NOTE] Could not update {path}: {e}")

if __name__ == "__main__":
    for p in PPT_FILES:
        apply_clean_layout(p)
