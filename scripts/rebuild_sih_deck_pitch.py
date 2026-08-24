"""
=============================================================================
BharatSRM-Net: Executive Business Pitch + Smart Education Alignment (SIH 2026)
=============================================================================
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_FILES = [
    r"C:\Users\muthu\Downloads\SIH2026-BharatSRM-Net-SmartEducation.pptx",
    r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx",
    r"C:\Users\muthu\Downloads\SIH2026-BharatSRM-Net-PitchDeck.pptx",
]
DIAG_DIR = "outputs/ppt_diagrams"

def style_bullet(p, title, desc, font_size=10.5, space_after=5, bold_color=RGBColor(15, 23, 42), text_color=RGBColor(51, 65, 85)):
    p.text = ""
    p.space_after = Pt(space_after)
    p.font.name = "Segoe UI"
    
    r1 = p.add_run()
    r1.text = f"• {title}: "
    r1.font.bold = True
    r1.font.size = Pt(font_size)
    r1.font.name = "Segoe UI"
    r1.font.color.rgb = bold_color
    
    r2 = p.add_run()
    r2.text = desc
    r2.font.bold = False
    r2.font.size = Pt(font_size)
    r2.font.name = "Segoe UI"
    r2.font.color.rgb = text_color

def build_smart_education_pitch(out_path):
    template_path = r"C:\Users\muthu\Downloads\SIH2026-IDEA-Presentation-Format (2).pptx"
    if not os.path.exists(template_path):
        return
    try:
        prs = pptx.Presentation(template_path)
        print(f"\nWeaving Smart Education + Business Pitch for: {out_path}...")

        # -------------------------------------------------------------
        # SLIDE 1: Title Slide (Smart Education Theme)
        # -------------------------------------------------------------
        s1 = prs.slides[0]
        for shape in s1.shapes:
            if shape.has_text_frame:
                if "SMART INDIA HACKATHON" in shape.text:
                    shape.top = Inches(0.35)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.7)
                    shape.text_frame.paragraphs[0].font.size = Pt(22)
                elif "TITLE PAGE" in shape.text or "BharatSRM" in shape.text:
                    shape.top = Inches(0.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(9.0)
                    shape.height = Inches(0.9)
                    shape.text_frame.clear()
                    p1 = shape.text_frame.paragraphs[0]
                    p1.text = "BharatSRM-Net"
                    p1.font.size = Pt(26)
                    p1.font.bold = True
                    p1.font.color.rgb = RGBColor(15, 23, 42)
                    p2 = shape.text_frame.add_paragraph()
                    p2.text = "AI Geospatial Platform Democratizing High-Resolution Satellite Intelligence for Space Education & National Defense"
                    p2.font.size = Pt(11.5)
                    p2.font.color.rgb = RGBColor(71, 85, 105)
                elif "Problem Statement ID" in shape.text:
                    shape.top = Inches(1.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(6.8)
                    shape.height = Inches(4.5)
                    tf = shape.text_frame
                    tf.clear()
                    fields = [
                        ("Problem Statement ID", "26142"),
                        ("Problem Statement Title", "Deep Learning Based Super Resolution Mapping (SRM) from Medium Resolution Satellite Imageries"),
                        ("Theme", "Smart Education (Space Science & Geospatial Learning)"),
                        ("Category", "Software"),
                        ("Organization / Ministry", "National Technical Research Organisation (NTRO)"),
                        ("Team Name", "ChaosToCode"),
                        ("Team ID", "[To be filled by team]"),
                    ]
                    for idx, (k, v) in enumerate(fields):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=11, space_after=3, bold_color=RGBColor(2, 132, 199))

        # -------------------------------------------------------------
        # SLIDE 2: Proposed Solution (Business Pitch + Smart Education)
        # -------------------------------------------------------------
        s2 = prs.slides[1]
        for s in list(s2.shapes):
            if s.name.startswith("Added_Diagram"):
                s2.shapes._spTree.remove(s._element)

        for shape in s2.shapes:
            if shape.has_text_frame:
                if "PROPOSED SOLUTION" in shape.text or "IDEA TITLE" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                    shape.text_frame.paragraphs[0].text = "PROPOSED SOLUTION: BHARATSRM-NET"
                elif "The Core Challenge" in shape.text or "The Multi-Crore" in shape.text or "Proposed Solution" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(5.3)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("The Cost & Education Barrier", "Commercial high-res satellite data costs crores (> ₹12 Lakh/1,000 km²), locking Indian university students, researchers, and government analysts out of fine-scale geospatial learning."),
                        ("Our Solution (BharatSRM-Net)", "Transforms free 10m public satellite data into 2.5m commercial-grade imagery with 16x sharper clarity — unlocking sovereign spatial intelligence at ₹0 extra cost."),
                        ("Smart Education Sandbox", "Serves as an interactive, hands-on learning lab where students and analysts learn Explainable AI, multi-spectral physics, and satellite image interpretation."),
                        ("Trust & Anti-Hallucination", "Built-in AI Confidence Meter (Uncertainty Map) teaches students and defense analysts exactly which features are 100% verified vs AI-inferred."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=6)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag2 = f"{DIAG_DIR}/diagram_slide2.png"
        if os.path.exists(diag2):
            pic = s2.shapes.add_picture(diag2, Inches(6.0), Inches(1.45), width=Inches(6.8), height=Inches(4.8))
            pic.name = "Added_Diagram_2"

        # -------------------------------------------------------------
        # SLIDE 3: Product Workflow & Technical Architecture
        # -------------------------------------------------------------
        s3 = prs.slides[2]
        for s in list(s3.shapes):
            if s.name.startswith("Added_Diagram"):
                s3.shapes._spTree.remove(s._element)

        for shape in s3.shapes:
            if shape.has_text_frame:
                if "TECHNICAL APPROACH" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.7)
                    shape.text_frame.paragraphs[0].text = "PRODUCT WORKFLOW & TECHNICAL ARCHITECTURE"
                elif "Multi-Modal Ingestion" in shape.text or "Technologies to be used" in shape.text or "End-to-End" in shape.text:
                    shape.top = Inches(0.95)
                    shape.left = Inches(0.5)
                    shape.width = Inches(12.3)
                    shape.height = Inches(1.0)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("End-to-End Automated Pipeline", "Ingests free 10m satellite imagery, combines it with ISRO elevation priors, applies an AI super-resolution engine, and outputs 2.5m high-res data with 4 intelligence layers in sub-second speed."),
                        ("Interactive Educational GIS Studio", "Real-time split-screen interface enabling students, researchers, and intelligence trainees to interactively compare raw vs AI-enhanced satellite layers."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=2)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag3 = f"{DIAG_DIR}/diagram_slide3.png"
        if os.path.exists(diag3):
            pic = s3.shapes.add_picture(diag3, Inches(0.5), Inches(2.1), width=Inches(12.3), height=Inches(4.6))
            pic.name = "Added_Diagram_3"

        # -------------------------------------------------------------
        # SLIDE 4: Feasibility, Deployment & Scalability
        # -------------------------------------------------------------
        s4 = prs.slides[3]
        for s in list(s4.shapes):
            if s.name.startswith("Added_Diagram"):
                s4.shapes._spTree.remove(s._element)

        for shape in s4.shapes:
            if shape.has_text_frame:
                if "FEASIBILITY" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                    shape.text_frame.paragraphs[0].text = "FEASIBILITY, VIABILITY & DEPLOYMENT"
                elif "Verified Working" in shape.text or "Analysis of the feasibility" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(5.3)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Verified Working Prototype", "Full-stack interactive Web GIS platform fully built and tested; verified with sub-second response times across 4 Indian biomes."),
                        ("100% Free Data Pipeline", "Uses freely accessible Copernicus Sentinel-2 and ISRO Bhuvan elevation data — zero recurring software API costs for schools or defense."),
                        ("Dual Deployment (Classroom to Defense)", "Deployable on standard university laptops / cloud labs as well as air-gapped secure defense command servers."),
                        ("Automated Weather Resilience", "Proprietary cloud-masking algorithms automatically clean atmospheric haze and cloud gaps to recover true terrain details."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=6)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag4 = f"{DIAG_DIR}/diagram_slide4.png"
        if os.path.exists(diag4):
            pic = s4.shapes.add_picture(diag4, Inches(6.0), Inches(1.45), width=Inches(6.8), height=Inches(4.8))
            pic.name = "Added_Diagram_4"

        # -------------------------------------------------------------
        # SLIDE 5: Business Impact, National Value & Smart Education
        # -------------------------------------------------------------
        s5 = prs.slides[4]
        for s in list(s5.shapes):
            if s.name.startswith("Added_Diagram"):
                s5.shapes._spTree.remove(s._element)

        for shape in s5.shapes:
            if shape.has_text_frame:
                if "IMPACT AND BENEFITS" in shape.text or "BUSINESS IMPACT" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                    shape.text_frame.paragraphs[0].text = "NATIONAL IMPACT, SMART EDUCATION & ROI"
                elif "Massive Cost Savings" in shape.text or "Potential impact" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.5)
                    shape.width = Inches(4.9)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Smart Education & Skilling (NEP 2020)", "Democratizes Space Science & GIS education across 500+ Indian universities, training the next generation of geospatial & defense data scientists."),
                        ("Massive Government ROI", "Saves ₹100s of Crores annually for ministries by replacing expensive foreign commercial satellite procurement ($15/km²)."),
                        ("PM Gati Shakti & Rural Infrastructure", "Automated extraction and vectorization of unpaved village roads and connectivity corridors (PMGSY)."),
                        ("Agriculture & Disaster Relief", "Precision farm parcel boundary tracking (PM Fasal Bima) + rapid flood inundation damage mapping for emergency teams."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=10.5, space_after=6)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        diag5 = f"{DIAG_DIR}/diagram_slide5.png"
        if os.path.exists(diag5):
            pic = s5.shapes.add_picture(diag5, Inches(5.6), Inches(1.4), width=Inches(7.2), height=Inches(5.0))
            pic.name = "Added_Diagram_5"

        # -------------------------------------------------------------
        # SLIDE 6: Research, Datasets & Validation
        # -------------------------------------------------------------
        s6 = prs.slides[5]
        for shape in s6.shapes:
            if shape.has_text_frame:
                if "RESEARCH" in shape.text or "VALIDATION" in shape.text:
                    shape.top = Inches(0.2)
                    shape.left = Inches(0.5)
                    shape.width = Inches(10.0)
                    shape.height = Inches(0.8)
                    shape.text_frame.paragraphs[0].text = "VALIDATION, RESEARCH & BENCHMARKS"
                elif "Rigorous Dataset" in shape.text or "Bayesian Uncertainty" in shape.text:
                    shape.top = Inches(1.25)
                    shape.left = Inches(0.6)
                    shape.width = Inches(12.0)
                    shape.height = Inches(5.2)
                    tf = shape.text_frame
                    tf.clear()
                    bullets = [
                        ("Rigorous Academic Pretraining", "Trained on 3,900+ real high-resolution satellite scene pairs covering agricultural, forest, urban, and desert biomes."),
                        ("ISRO & Government Compliance", "Conforms to official ISRO NRSC Bhuvan Land-Use Land-Cover classification schemas and Copernicus Level-2A standards."),
                        ("Peer-Reviewed Scientific Foundation", "Built on proven research in Bayesian Uncertainty (NeurIPS), Cloud Inpainting (ECCV), and Earth Observation Super-Resolution (NeurIPS)."),
                        ("Open Public Data Sources", "European Space Agency (Copernicus Open Access Hub), ISRO NRSC Bhuvan (CartoDEM), SPOT 6/7 1.5m High-Resolution References."),
                    ]
                    for idx, (k, v) in enumerate(bullets):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        style_bullet(p, k, v, font_size=11.5, space_after=8)
                elif "ChaosToCode" in shape.text or "Your Team Name" in shape.text:
                    shape.text_frame.paragraphs[0].text = "ChaosToCode"

        prs.save(out_path)
        print(f"[SUCCESS] Smart Education pitch deck saved to: {out_path}")
    except Exception as e:
        print(f"[NOTE] Could not update {out_path}: {e}")

if __name__ == "__main__":
    for p in PPT_FILES:
        build_smart_education_pitch(p)
